from pptx import Presentation
import matplotlib.pyplot as plt
from pptx.util import Inches, Pt
import os

        #### PowerPoint ####

        
class PowerPoint():
    def __init__(self, data_filtered, summary_qty, summary_code, summary_step, summary_product, summary_data, sorted_df, df_top, df_f):
        #Create a new PowerPoint presentation
        self.data_filtered = data_filtered
        self.summary_qty = summary_qty
        self.summary_code = summary_code
        self.summary_step = summary_step
        self.summary_product = summary_product
        self.summary_data = summary_data
        self.sorted_df = sorted_df
        self.df_top = df_top
        self.df_f = df_f
        self.prs = self.Create_pp()


    def Create_pp(self):

        #Create fresh PowerPoint presentation
        prs = Presentation()

        #All data we are going use in this powerpoint report
        slide1_data_report = self.summary_qty.reset_index()
        slide12_data_report = self.summary_code.reset_index()
        slide3_data_report = self.summary_step.reset_index()
        slide13_data_report = self.summary_product.reset_index()
        slide4_data_report = self.data_filtered[["LOT_ID", "QTY", "COMMENTS", "CODE"]]

        labels = self.sorted_df["CODES"]
        values = self.sorted_df["QTY"]

        ### Make Barplot ###

        bars = plt.bar(labels, values, orientation='vertical')
        #plt.figure(figsize=(10,6))
        plt.xticks(rotation=40, fontsize= 7)
        plt.ylabel('QTY')
        title_string = "Glassing scrap report YEAR-WEEK: {}".format(str(self.data_filtered["WEEK"][0]))
        plt.title(title_string)

        #add values on top of the bars
        for bar in bars:
            height = bar.get_height()
            plt.annotate(
                f'{height}',
                xy=(bar.get_x() + bar.get_width() / 2, height),
                xytext = (0, 3),
                textcoords='offset points',
                ha='center', va='bottom',
                fontsize=8
                )

        #savebarplot for loading it later to slide
        bar_plot_filename = 'bar_plot.png'
        plt.savefig(bar_plot_filename, bbox_inches='tight')
        plt.close


        ### SLIDE 0 ###

        #Create first slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)

        #Add text to first slide
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Scrap Report"
        subtitle.text = "Glassing"

        ## SLIDE 1 ##
        # barplot #

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        left = Inches(0.5)
        top = Inches(0.5)
        slide.shapes.add_picture(bar_plot_filename, left, top, width=Inches(8))
        
        ### SLIDE 2 ###

        #Create second slide
        slide = prs.slides.add_slide(prs.slide_layouts[3])

        #Add text
        title = slide.shapes.title
        title.text = "Scrap Report"

        content1_placeholder = slide.placeholders[1]
        content2_placeholder = slide.placeholders[2]

        content1_placeholder.text = "Device"
        content2_placeholder.text = "Scrap Code"

        #table shape and position in slide
        rows, cols = slide1_data_report.shape
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(1.2)
        table1 = slide.shapes.add_table(rows+1, cols, left, top, width, height).table

        rows, cols = slide1_data_report.shape
        left = Inches(5)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(1.2)
        table12 = slide.shapes.add_table(rows+1, cols, left, top, width, height).table

        #Set width of each column separately
        table1.columns[0].width = Inches(2)
        table1.columns[1].width = Inches(0.7)
        table12.columns[0].width = Inches(2)
        table12.columns[1].width = Inches(0.7)


        # Add the headers to table
        for i, col_name in enumerate(slide1_data_report.columns):
            table1.cell(0, i).text = col_name
            
        # Add the data to the table
        for i, row in enumerate(slide1_data_report.itertuples(index=False)):
            for j, value in enumerate(row):
                table1.cell(i+1, j).text = str(value)
                
        # Add the headers to table
        for i, col_name in enumerate(slide12_data_report.columns):
            table12.cell(0, i).text = col_name
            
        # Add the data to the table
        for i, row in enumerate(slide12_data_report.itertuples(index=False)):
            for j, value in enumerate(row):
                table12.cell(i+1, j).text = str(value)
                
        
        ### SLIDE 3 ###

        #Create second slide
        slide = prs.slides.add_slide(prs.slide_layouts[3])

        #Add text
        title = slide.shapes.title
        title.text = "Scrap Report"

        content1_placeholder = slide.placeholders[1]
        content2_placeholder = slide.placeholders[2]

        content1_placeholder.text = "Current Step"
        content2_placeholder.text = "Product"

        #table shape and position in slide
        rows, cols = slide3_data_report.shape
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(1.2)
        table2 = slide.shapes.add_table(rows+1, cols, left, top, width, height).table

        rows, cols = slide3_data_report.shape
        left = Inches(5.5)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(1.2)
        table22 = slide.shapes.add_table(rows+1, cols, left, top, width, height).table

        #Set width of each column separately
        table2.columns[0].width = Inches(3.5)
        table2.columns[1].width = Inches(0.7)
        table22.columns[0].width = Inches(2)
        table22.columns[1].width = Inches(0.7)


        # Add the headers to table
        for i, col_name in enumerate(slide3_data_report.columns):
            table2.cell(0, i).text = col_name
            
        # Add the data to the table
        for i, row in enumerate(slide3_data_report.itertuples(index=False)):
            for j, value in enumerate(row):
                table2.cell(i+1, j).text = str(value)
                
        # Add the headers to table
        for i, col_name in enumerate(slide13_data_report.columns):
            table22.cell(0, i).text = col_name
            
        # Add the data to the table
        for i, row in enumerate(slide13_data_report.itertuples(index=False)):
            for j, value in enumerate(row):
                table22.cell(i+1, j).text = str(value)
                
        

        ### SLIDE 4 ###

        #Create third slide with table
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        rows, cols = slide4_data_report.shape
        left = Inches(0.4)
        top = Inches(0.2)
        width = Inches(8.5)
        height = Inches(1.2)
        table3 = slide.shapes.add_table(rows+1, cols, left, top, width, height).table

        #Set width of each column separately
        table3.columns[0].width = Inches(1.5)
        table3.columns[1].width = Inches(0.7)
        table3.columns[2].width = Inches(6)
        table3.columns[3].width = Inches(1)


        # Add the headers to table
        for i, col_name in enumerate(slide4_data_report.columns):
            table3.cell(0, i).text = col_name
            
        # Add the data to the table
        for i, row in enumerate(slide4_data_report.itertuples(index=False)):
            for j, value in enumerate(row):
                table3.cell(i+1, j).text = str(value)


        ### SLIDE 0 for HOLDS ###

        #Create first slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)

        #Add text to first slide
        title2 = slide.shapes.title
        subtitle2 = slide.placeholders[1]
        title2.text = "Hold Report"
        subtitle2.text = "Glassing"
        
        #PowerPoint slides for data
        slide5_data_report = self.df_f
        slide51_data_report = self.df_top
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        rows, cols = slide51_data_report.shape
        left = Inches(0.4)
        top = Inches(0.2)
        width = Inches(8.5)
        height = Inches(1.2)
        table41 = slide.shapes.add_table(rows+1, cols, left, top, width, height).table
        
        # Add the headers to table
        for i, col_name in enumerate(slide51_data_report.columns):
            table41.cell(0, i).text = col_name
        
        # Add the data to the table
        for i, row in enumerate(slide51_data_report.itertuples(index=False)):
            for j, value in enumerate(row):
                table41.cell(i+1, j).text = str(value)
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        rows2, cols2 = slide5_data_report.shape
        left = Inches(0.4)
        top = Inches(0.2)
        width = Inches(8.5)
        height = Inches(1.2)
        table4 = slide.shapes.add_table(rows2+1, cols2, left, top, width, height).table
        
        # Add the headers to table
        for i, col_name in enumerate(slide5_data_report.columns):
            table4.cell(0, i).text = col_name
        
        # Add the data to the table
        for i, row in enumerate(slide5_data_report.itertuples(index=False)):
            for j, value in enumerate(row):
                table4.cell(i+1, j).text = str(value)
        
        #Set width of each column separately
        table41.columns[0].width = Inches(7)
        table41.columns[1].width = Inches(1.2)
        table4.columns[0].width = Inches(7)
        table4.columns[1].width = Inches(1.2)
        
        return prs

    def Ask_Save_PP(self):
        save_to = input("Do you want to save the data to a PowerPoint slideshow Y/N:")
        save_name = "Weekly report - Glassing.pptx"

        if save_to.lower() == "yes" or save_to.lower() == "y":
            self.prs.save(save_name)
            print("\n**Data saved to {}**".format(save_name))
            os.remove("bar_plot.png")
            os.startfile(save_name)
        else:
            are_you_sure = input("Are you sure that you dont want to save PowerPoint Y/N:")
            if are_you_sure.lower() == "yes" or are_you_sure.lower() == "y":
                print("\n**Data is not saved**")
            else:
                self.prs.save(save_name)
                print("\n**Data saved to {}**".format(save_name))
                os.remove("bar_plot.png")
                os.startfile(save_name)